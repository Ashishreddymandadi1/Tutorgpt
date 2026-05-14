"""
Parses PDF, PPTX, and DOCX files into a list of (page_num, text) tuples.
Page numbers are 1-indexed.
"""
import os
from typing import List, Tuple


def parse_document(file_path: str) -> List[Tuple[int, str]]:
    """Returns list of (page_num, text) tuples."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".docx":
        return _parse_docx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _parse_pdf(file_path: str) -> List[Tuple[int, str]]:
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append((i, text))
    return pages


def _parse_pptx(file_path: str) -> List[Tuple[int, str]]:
    from pptx import Presentation
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text)
                    if line.strip():
                        texts.append(line)
        if texts:
            pages.append((i, "\n".join(texts)))
    return pages


def _parse_docx(file_path: str) -> List[Tuple[int, str]]:
    from docx import Document
    doc = Document(file_path)
    # DOCX has no true page boundaries; group ~30 paragraphs per "page"
    page_size = 30
    all_paras = [p.text for p in doc.paragraphs if p.text.strip()]
    pages = []
    for i in range(0, len(all_paras), page_size):
        chunk_paras = all_paras[i:i + page_size]
        pages.append((i // page_size + 1, "\n".join(chunk_paras)))
    return pages
